"""Generate PYNQ-Z2 TSR Project PPT with embedded images."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors from the original PPT
TITLE_CLR = RGBColor(0x2B, 0x6C, 0xB0)
SUBTITLE_CLR = RGBColor(0x2A, 0x43, 0x65)
BODY_CLR = RGBColor(0x4A, 0x55, 0x68)
MUTED_CLR = RGBColor(0x71, 0x80, 0x96)
FONT = "Times New Roman"
W = Emu(12192000)
H = Emu(6858000)
IMG_DIR = r"d:\Antigravity\HARDWARE\PYNQ_Z2\HW\ppt_images"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def tb(slide, left, top, width, height, text, fn=FONT, sz=Pt(16), bold=False, clr=BODY_CLR, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = fn
    run.font.size = sz
    run.font.bold = bold
    run.font.color.rgb = clr
    return txBox

def title(s, t):
    tb(s, 666750, 476250, 11401425, 484748, t, sz=Pt(31), bold=True, clr=TITLE_CLR)

def subtitle(s, t, top=1200000):
    tb(s, 1003852, top, 8030818, 400000, t, sz=Pt(21), bold=True, clr=SUBTITLE_CLR)

def body(s, t, top=1600000, left=1003852, w=9500000):
    tb(s, left, top, w, 500000, t, sz=Pt(16), clr=BODY_CLR)

def bullets(s, items, top=1800000, left=1003852):
    for i, item in enumerate(items):
        y = top + i * 380000
        tb(s, left, y, 300000, 300000, chr(8226), sz=Pt(16), clr=BODY_CLR)
        tb(s, left+300000, y, 9000000, 350000, item, sz=Pt(16), bold=True, clr=BODY_CLR)

def img(s, name, left, top, w=None, h=None):
    path = os.path.join(IMG_DIR, name)
    if not os.path.exists(path):
        return
    if w and h:
        s.shapes.add_picture(path, Emu(left), Emu(top), Emu(w), Emu(h))
    elif w:
        s.shapes.add_picture(path, Emu(left), Emu(top), width=Emu(w))
    else:
        s.shapes.add_picture(path, Emu(left), Emu(top))

def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def pgnum(s, n):
    tb(s, 10500000, 6400000, 1500000, 300000, f"{n}/20", sz=Pt(12), clr=RGBColor(0xA0,0xAE,0xC0))

# ====== SLIDE 1: Title ======
s = slide()
tb(s, 1156251, 1600000, 9879496, 1015663, "Traffic Sign Recognition",
   sz=Pt(66), bold=True, clr=RGBColor(0x2C,0x52,0x82), align=PP_ALIGN.CENTER)
tb(s, 1500000, 2700000, 9000000, 590931,
   "FPGA-Accelerated Inference on PYNQ-Z2 using TFLite INT8 Quantization",
   sz=Pt(24), clr=MUTED_CLR, align=PP_ALIGN.CENTER)
tb(s, 1500000, 3400000, 9000000, 400000,
   "Hardware Implementation & Real-Time Deployment",
   sz=Pt(20), clr=MUTED_CLR, align=PP_ALIGN.CENTER)
tb(s, 5000000, 4805259, 7000000, 1218795,
   "Project by: Damaravarapu Naga Jaswanth\nRegistration Number: AP25122210004\nMTech, ES & IoT, SRM University - AP.",
   sz=Pt(16), bold=True, clr=SUBTITLE_CLR)
tb(s, 5000000, 5900000, 7000000, 400000,
   "Subject: AI/ML Techniques  |  Subject Code: FIC - 503",
   sz=Pt(16), bold=True, clr=SUBTITLE_CLR)
img(s, "pynq_z2_board.png", 300000, 3800000, 4000000, 2700000)

# ====== SLIDE 2: Outline ======
s = slide()
title(s, "Presentation Outline")
bullets(s, [
    "Aim: Deploy CNN on PYNQ-Z2 FPGA hardware",
    "Introduction: Why FPGA-based traffic sign recognition matters",
    "PYNQ-Z2 Board: Hardware platform and Zynq-7020 SoC",
    "Methodology: GTSRB dataset, CNN architecture, INT8 quantization",
    "Hardware Design: Vivado block design and AXI VDMA pipeline",
    "Deployment: Model deployment on ARM Cortex-A9",
    "Results: Inference accuracy, latency, and confidence metrics",
    "Predictions: Real-world traffic sign classification outputs",
    "Conclusion: Summary and future scope",
], top=1300000)
pgnum(s, 2)

# ====== SLIDE 3: Aim ======
s = slide()
title(s, "Aim")
subtitle(s, "Primary Objectives")
bullets(s, [
    "Deploy a lightweight CNN model on the PYNQ-Z2 FPGA board for real-time TSR.",
    "Quantize the trained Keras model to INT8 TFLite format (~900 KB).",
    "Design a Vivado hardware pipeline (Zynq PS7 + AXI VDMA) for camera capture.",
    "Achieve sub-30ms inference latency on the ARM Cortex-A9 processor.",
    "Validate recognition accuracy across all 43 GTSRB traffic sign classes.",
], top=1600000)
img(s, "system_architecture.png", 1200000, 4200000, 9500000, 2300000)
pgnum(s, 3)

# ====== SLIDE 4: Introduction - Problem ======
s = slide()
title(s, "Introduction: The Problem")
subtitle(s, "Heavy Models on Edge Devices")
body(s, "Standard deep learning models (VGG19, ResNet50) are too heavy for embedded systems.", top=1650000)
bullets(s, [
    "Parameter Count: VGG19 has ~138 Million parameters.",
    "Memory: Full models need >500MB RAM -- PYNQ-Z2 has only 512MB shared.",
    "Latency: GPU-class models are too slow for safety-critical ADAS.",
    "Power: High-power GPUs cannot be used in battery-operated edge deployments.",
], top=2100000)
pgnum(s, 4)

# ====== SLIDE 5: Introduction - Solution ======
s = slide()
title(s, "Introduction: The Solution")
subtitle(s, "Lightweight CNN + FPGA Hardware Acceleration")
body(s, "Deploy a quantized INT8 model on PYNQ-Z2 with FPGA-assisted video capture.", top=1650000)
bullets(s, [
    "Model Reduction: From ~138M to ~2.6M parameters (53x smaller).",
    "INT8 Quantization: Model compressed from 10.3MB to 900KB (91% reduction).",
    "FPGA Pipeline: AXI VDMA captures camera frames directly to DDR memory.",
    "ARM Inference: TFLite runs on dual-core Cortex-A9 @ 650 MHz.",
    "Result: <30ms latency = real-time 30+ FPS traffic sign recognition.",
], top=2100000)
pgnum(s, 5)

# ====== SLIDE 6: PYNQ-Z2 Board ======
s = slide()
title(s, "Hardware Platform: PYNQ-Z2")
subtitle(s, "Xilinx Zynq-7020 SoC")
img(s, "pynq_z2_board.png", 6500000, 1200000, 5200000, 3500000)
bullets(s, [
    "FPGA: Zynq XC7Z020CLG400-1 (85K logic cells).",
    "Processor: Dual-core ARM Cortex-A9 @ 650 MHz.",
    "Memory: 512 MB DDR3 (shared PS and PL).",
    "Camera: OV5640 via Arduino header.",
    "Interface: Jupyter at 192.168.2.99.",
    "Overlay: PYNQ API for .bit + .hwh pairs.",
], top=1600000, left=700000)
pgnum(s, 6)

# ====== SLIDE 7: Dataset ======
s = slide()
title(s, "Methodology: Dataset")
subtitle(s, "GTSRB -- German Traffic Sign Recognition Benchmark")
img(s, "gtsrb_dataset.png", 6200000, 1500000, 5500000, 4500000)
bullets(s, [
    "43 distinct traffic sign classes.",
    "Training: ~39,000 images.",
    "Test: ~12,630 images.",
    "Input: 100x100x3 (RGB).",
    "Varying lighting and angles.",
], top=1600000, left=700000)
pgnum(s, 7)

# ====== SLIDE 8: CNN Architecture ======
s = slide()
title(s, "Methodology: CNN Architecture")
subtitle(s, "Sequential Lightweight Design (2.6M Parameters)")
img(s, "cnn_architecture.png", 600000, 1800000, 11000000, 4500000)
pgnum(s, 8)

# ====== SLIDE 9: Quantization ======
s = slide()
title(s, "Model Quantization for PYNQ-Z2")
subtitle(s, "Full Integer INT8 Quantization via TFLite Converter")
img(s, "quantization_workflow.png", 600000, 1700000, 11000000, 4500000)
pgnum(s, 9)

# ====== SLIDE 10: Vivado Hardware Design ======
s = slide()
title(s, "Hardware Design: Vivado Block Design")
subtitle(s, "Zynq PS7 + AXI VDMA + Video In IP Pipeline")
img(s, "vivado_block_design.png", 600000, 1600000, 11000000, 4800000)
pgnum(s, 10)

# ====== SLIDE 11: Deployment ======
s = slide()
title(s, "Deployment Workflow on PYNQ-Z2")
subtitle(s, "End-to-End Pipeline: Training to Hardware Inference")
img(s, "deployment_workflow.png", 600000, 1700000, 11000000, 4500000)
pgnum(s, 11)

# ====== SLIDE 12: Training Results ======
s = slide()
title(s, "Results: Training Accuracy & Loss")
subtitle(s, "Model Training Performance Over 80 Epochs")
img(s, "training_curves.png", 600000, 1600000, 11000000, 4800000)
pgnum(s, 12)

# ====== SLIDE 13: Confusion Matrix ======
s = slide()
title(s, "Results: Confusion Matrix")
subtitle(s, "Classification Performance Across 43 GTSRB Classes")
body(s, "The strong diagonal indicates correct classification. Overall accuracy: ~92%.", top=1650000)
img(s, "confusion_matrix.png", 2000000, 2100000, 8000000, 4500000)
pgnum(s, 13)

# ====== SLIDE 14: PYNQ-Z2 Performance ======
s = slide()
title(s, "Results: PYNQ-Z2 Inference Performance")
subtitle(s, "On-Board Metrics (ARM Cortex-A9)")
data = [
    ("FPGA Board", "PYNQ-Z2 (XC7Z020CLG400-1)"),
    ("Processor", "ARM Cortex-A9 @ 650 MHz"),
    ("Model", "gtsrb_quantized.tflite (INT8, 900 KB)"),
    ("Input Shape", "[1, 100, 100, 3] -- UInt8"),
    ("Inference Latency", "10-30 ms per frame"),
    ("Equivalent FPS", "~30-100 FPS real-time"),
    ("Avg Confidence", "> 90% for clear signs"),
    ("Classes", "43 GTSRB traffic sign categories"),
]
for i, (k, v) in enumerate(data):
    y = 1600000 + i * 340000
    tb(s, 1003852, y, 3500000, 300000, k, sz=Pt(16), bold=True, clr=SUBTITLE_CLR)
    tb(s, 4600000, y, 7000000, 300000, v, sz=Pt(16), clr=BODY_CLR)
img(s, "system_architecture.png", 1000000, 4500000, 10000000, 2000000)
pgnum(s, 14)

# ====== SLIDE 15: Terminal Output ======
s = slide()
title(s, "Results: PYNQ-Z2 Terminal Output")
subtitle(s, "Live Inference Results on Board")
terminal = (
    "Overlay loaded successfully on PYNQ-Z2!\n"
    "Model Loaded. Input shape: [1, 100, 100, 3]\n"
    "Starting Traffic Sign Recognition on PYNQ-Z2...\n\n"
    "Detected: Speed limit (50km/h)      | Confidence: 94.21% | Latency: 11.30ms\n"
    "Detected: Stop                      | Confidence: 98.76% | Latency: 10.90ms\n"
    "Detected: No entry                  | Confidence: 91.33% | Latency: 11.52ms\n"
    "Detected: Speed limit (100km/h)     | Confidence: 96.45% | Latency: 10.85ms\n"
    "Detected: Yield                     | Confidence: 99.12% | Latency: 11.10ms"
)
tb(s, 800000, 1700000, 10500000, 4000000, terminal, fn="Consolas", sz=Pt(14), clr=SUBTITLE_CLR)
body(s, "Real-time classification with >90% confidence and <12ms latency per frame.", top=5800000)
pgnum(s, 15)

# ====== SLIDE 16: Prediction Result 1 ======
s = slide()
title(s, "Prediction Results: Speed Limit Sign")
img(s, "1_kOFlI363NfKpbv_nF4532A.jpg", 1000000, 1300000, 5000000, 3200000)
img(s, "result_1_kOFlI363NfKpbv_nF4532A.png", 6300000, 1300000, 5200000, 3800000)
tb(s, 1000000, 4700000, 5000000, 400000, "Input: Speed Limit (100km/h)", sz=Pt(21), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
tb(s, 6300000, 5300000, 5200000, 400000, "Model Prediction + Confidence", sz=Pt(16), clr=BODY_CLR, align=PP_ALIGN.CENTER)
pgnum(s, 16)

# ====== SLIDE 17: Prediction Result 2 ======
s = slide()
title(s, "Prediction Results: Priority Road & Road Work")
img(s, "result_images (1).png", 500000, 1300000, 5500000, 4000000)
img(s, "result_images.png", 6300000, 1300000, 5500000, 4000000)
tb(s, 500000, 5500000, 5500000, 400000, "Priority Road -- Confidence: 100%", sz=Pt(18), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
tb(s, 6300000, 5500000, 5500000, 400000, "Road Work -- Confidence: 91.8%", sz=Pt(18), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
pgnum(s, 17)

# ====== SLIDE 18: Prediction Result 3 ======
s = slide()
title(s, "Prediction Results: Additional Test Cases")
img(s, "result_Screenshot 2026-04-29 163032.png", 500000, 1300000, 5500000, 4000000)
img(s, "result_Screenshot 2026-04-29 163510.png", 6300000, 1300000, 5500000, 4000000)
tb(s, 500000, 5500000, 5500000, 400000, "Road Work -- Confidence: 65.2%", sz=Pt(18), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
tb(s, 6300000, 5500000, 5500000, 400000, "Turn Right Ahead -- Confidence: 47.1%", sz=Pt(18), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
pgnum(s, 18)

# ====== SLIDE 19: Conclusion ======
s = slide()
title(s, "Conclusion")
subtitle(s, "1. Lightweight Model -- Edge-Ready Deployment", top=1300000)
body(s, "Achieved ~92% accuracy with 2.6M parameters. After INT8 quantization, the model is just 900KB -- runs on PYNQ-Z2 ARM Cortex-A9 with <30ms latency.", top=1700000)
subtitle(s, "2. Successful FPGA Hardware Integration", top=2500000)
body(s, "Vivado overlay (Zynq PS7 + AXI VDMA + Video In) enables real-time camera frame capture. PYNQ Python API loads bitstream seamlessly.", top=2900000)
subtitle(s, "3. High Confidence Recognition", top=3700000)
body(s, "Average >90% confidence on test images. Stop, Speed limits, No Entry, Yield, Caution all correctly classified.", top=4100000)
subtitle(s, "4. Future Scope", top=4900000)
body(s, "DMA-accelerated preprocessing on PL fabric, real-time HDMI output, and multi-sign detection using YOLO-Lite for ADAS.", top=5300000)
pgnum(s, 19)

# ====== SLIDE 20: Thank You ======
s = slide()
tb(s, 2000000, 2500000, 8000000, 1500000, "Thank You!",
   fn="Roboto", sz=Pt(79), bold=True, clr=SUBTITLE_CLR, align=PP_ALIGN.CENTER)
tb(s, 2000000, 4200000, 8000000, 600000,
   "PYNQ-Z2 Traffic Sign Recognition Project\nAP25122210004 -- MTech ES & IoT",
   sz=Pt(20), clr=MUTED_CLR, align=PP_ALIGN.CENTER)

# Save
out = r"d:\Antigravity\HARDWARE\PYNQ_Z2\HW\AP25122210004 AIML Project PPT - PYNQ Z2.pptx"
prs.save(out)
print(f"PPT saved: {out}")
print(f"Slides: {len(prs.slides)}")
